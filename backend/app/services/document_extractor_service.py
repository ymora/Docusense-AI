"""
Document extraction service for DocuSense AI
Handles text extraction from Office documents and other formats
"""

from pathlib import Path
import asyncio
from sqlalchemy.orm import Session

from .base_service import BaseService, log_service_operation
from ..core.types import ServiceResponse, TextExtractionResult
from ..core.media_formats import get_supported_formats_keys, is_format_supported_in_dict


class DocumentExtractorService(BaseService):
    """
    Service d'extraction de texte des documents Office
    """

    def __init__(self, db: Session):
        super().__init__(db)
        self.supported_formats = {
            # Documents Word
            'docx': self._extract_from_docx,
            'doc': self._extract_from_doc,
            'rtf': self._extract_from_rtf,
            'odt': self._extract_from_odt,
            
            # Tableurs Excel
            'xlsx': self._extract_from_xlsx,
            'xls': self._extract_from_xls,
            'ods': self._extract_from_ods,
            'csv': self._extract_from_csv,
            
            # Présentations PowerPoint
            'pptx': self._extract_from_pptx,
            'ppt': self._extract_from_ppt,
            'odp': self._extract_from_odp,
            
            # Autres formats
            'pdf': self._extract_from_pdf,
            'txt': self._extract_from_txt,
            'html': self._extract_from_html,
            'htm': self._extract_from_html,
        }

    @log_service_operation("extract_text")
    async def extract_text(self, file_path: str) -> ServiceResponse:
        """
        Extrait le texte d'un fichier selon son format
        
        Args:
            file_path: Chemin vers le fichier
            
        Returns:
            ServiceResponse: Texte extrait ou erreur
        """
        try:
            text = await self.safe_execute_async("extract_text_logic", self._extract_text_logic, file_path)
            return {"success": True, "data": {"text": text, "file_path": file_path}}
        except Exception as e:
            return {"success": False, "error": f"Erreur lors de l'extraction: {str(e)}"}

    async def _extract_text_logic(self, file_path: str) -> str:
        """Logic for extracting text from file"""
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"Fichier non trouvé: {file_path}")

        extension = file_path.suffix.lower().lstrip('.')
        
        if extension not in self.supported_formats:
            raise ValueError(f"Format non supporté: {extension}")

        # Extraire le texte de manière asynchrone
        extractor_func = self.supported_formats[extension]
        text = await asyncio.to_thread(extractor_func, file_path)
        
        if text:
            self.logger.info(f"Texte extrait de {file_path.name}: {len(text)} caractères")
            return text
        else:
            self.logger.warning(f"Aucun texte extrait de {file_path.name}")
            return ""

    def _extract_from_docx(self, file_path: Path) -> str:
        """Extrait le texte d'un fichier DOCX"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            text_parts = []
            
            # Extraire le texte des paragraphes
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # Extraire le texte des tableaux
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            
            return "\n".join(text_parts) if text_parts else ""
            
        except Exception as e:
            self.logger.error(f"Erreur extraction DOCX {file_path}: {str(e)}")
            return ""

    def _extract_from_doc(self, file_path: Path) -> str:
        """Extrait le texte d'un fichier DOC (Word 97-2003)"""
        try:
            # Pour les fichiers DOC, on utilise une approche simple
            # En production, vous pourriez utiliser python-docx2txt ou antiword
            import subprocess
            import tempfile
            
            # Essayer d'utiliser antiword si disponible
            try:
                result = subprocess.run(['antiword', str(file_path)], 
                                      capture_output=True, text=True, timeout=30)
                if result.returncode == 0 and result.stdout.strip():
                    return result.stdout.strip()
            except (subprocess.TimeoutExpired, FileNotFoundError):
                pass
            
            # Fallback: essayer de lire comme texte brut
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                    if content.strip():
                        return content.strip()
            except:
                pass
            
            return f"Document Word (format .doc): {file_path.name}"
            
        except Exception as e:
            self.logger.error(f"Erreur extraction DOC {file_path}: {str(e)}")
            return ""

    def _extract_from_xlsx(self, file_path: Path) -> str:
        """Extrait le texte d'un fichier XLSX"""
        try:
            from openpyxl import load_workbook
            
            workbook = load_workbook(file_path, data_only=True)
            text_parts = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = [f"Feuille: {sheet_name}"]
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell_value in row:
                        if cell_value is not None:
                            row_text.append(str(cell_value).strip())
                    if row_text:
                        sheet_text.append(" | ".join(row_text))
                
                if len(sheet_text) > 1:  # Plus que juste le nom de la feuille
                    text_parts.extend(sheet_text)
                    text_parts.append("")  # Ligne vide entre les feuilles
            
            return "\n".join(text_parts) if text_parts else ""
            
        except Exception as e:
            self.logger.error(f"Erreur extraction XLSX {file_path}: {str(e)}")
            return ""

    def _extract_from_xls(self, file_path: Path) -> str:
        """Extrait le texte d'un fichier XLS (Excel 97-2003)"""
        try:
            import xlrd
            
            workbook = xlrd.open_workbook(file_path)
            text_parts = []
            
            for sheet_name in workbook.sheet_names():
                sheet = workbook.sheet_by_name(sheet_name)
                sheet_text = [f"Feuille: {sheet_name}"]
                
                for row_idx in range(sheet.nrows):
                    row_text = []
                    for col_idx in range(sheet.ncols):
                        cell_value = sheet.cell_value(row_idx, col_idx)
                        if cell_value:
                            row_text.append(str(cell_value).strip())
                    if row_text:
                        sheet_text.append(" | ".join(row_text))
                
                if len(sheet_text) > 1:  # Plus que juste le nom de la feuille
                    text_parts.extend(sheet_text)
                    text_parts.append("")  # Ligne vide entre les feuilles
            
            return "\n".join(text_parts) if text_parts else ""
            
        except Exception as e:
            self.logger.error(f"Erreur extraction XLS {file_path}: {str(e)}")
            return ""

    def _extract_from_csv(self, file_path: Path) -> str:
        """Extrait le texte d'un fichier CSV"""
        try:
            import csv
            
            text_parts = []
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.reader(f)
                for row in reader:
                    if row:
                        row_text = " | ".join([cell.strip() for cell in row if cell.strip()])
                        if row_text:
                            text_parts.append(row_text)
            
            return "\n".join(text_parts) if text_parts else ""
            
        except Exception as e:
            self.logger.error(f"Erreur extraction CSV {file_path}: {str(e)}")
            return ""

    def _extract_from_pptx(self, file_path: Path) -> str:
        """Extrait le texte d'un fichier PPTX"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"Diapositive {slide_num}"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if len(slide_text) > 1:  # Plus que juste le numéro de diapositive
                    text_parts.extend(slide_text)
                    text_parts.append("")  # Ligne vide entre les diapositives
            
            return "\n".join(text_parts) if text_parts else ""
            
        except Exception as e:
            self.logger.error(f"Erreur extraction PPTX {file_path}: {str(e)}")
            return ""

    def _extract_from_ppt(self, file_path: Path) -> str:
        """Extrait le texte d'un fichier PPT (PowerPoint 97-2003)"""
        try:
            # Pour les fichiers PPT, on utilise une approche simple
            # En production, vous pourriez utiliser python-pptx ou d'autres outils
            return f"Présentation PowerPoint (format .ppt): {file_path.name}"
            
        except Exception as e:
            self.logger.error(f"Erreur extraction PPT {file_path}: {str(e)}")
            return ""

    def _extract_from_pdf(self, file_path: Path) -> str:
        """Extrait le texte d'un fichier PDF"""
        try:
            import PyPDF2
            
            text_parts = []
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_parts.append(f"Page {page_num}:")
                        text_parts.append(page_text.strip())
                        text_parts.append("")  # Ligne vide entre les pages
            
            return "\n".join(text_parts) if text_parts else ""
            
        except Exception as e:
            self.logger.error(f"Erreur extraction PDF {file_path}: {str(e)}")
            return ""

    def _extract_from_txt(self, file_path: Path) -> str:
        """Extrait le texte d'un fichier TXT"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
                return content.strip() if content.strip() else ""
                
        except Exception as e:
            self.logger.error(f"Erreur extraction TXT {file_path}: {str(e)}")
            return ""

    def _extract_from_html(self, file_path: Path) -> str:
        """Extrait le texte d'un fichier HTML"""
        try:
            import re
            
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                html_content = f.read()
            
            # Supprimer les balises HTML
            text = re.sub(r'<[^>]+>', ' ', html_content)
            
            # Supprimer les espaces multiples
            text = re.sub(r'\s+', ' ', text)
            
            # Supprimer les entités HTML communes
            text = text.replace('&nbsp;', ' ')
            text = text.replace('&amp;', '&')
            text = text.replace('&lt;', '<')
            text = text.replace('&gt;', '>')
            text = text.replace('&quot;', '"')
            text = text.replace('&#39;', "'")
            
            return text.strip() if text.strip() else ""
            
        except Exception as e:
            self.logger.error(f"Erreur extraction HTML {file_path}: {str(e)}")
            return ""

    def _extract_from_rtf(self, file_path: Path) -> str:
        """Extrait le texte d'un fichier RTF"""
        try:
            import subprocess
            
            # Essayer d'utiliser untext si disponible
            try:
                result = subprocess.run(['untext', str(file_path)], 
                                      capture_output=True, text=True, timeout=30)
                if result.returncode == 0 and result.stdout.strip():
                    return result.stdout.strip()
            except (subprocess.TimeoutExpired, FileNotFoundError):
                pass
            
            # Fallback: essayer de lire comme texte brut
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                    if content.strip():
                        return content.strip()
            except:
                pass
            
            return f"Document RTF: {file_path.name}"
            
        except Exception as e:
            self.logger.error(f"Erreur extraction RTF {file_path}: {str(e)}")
            return ""

    def _extract_from_odt(self, file_path: Path) -> str:
        """Extrait le texte d'un fichier ODT (OpenDocument Text)"""
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            
            # Les fichiers ODT sont des archives ZIP
            with zipfile.ZipFile(file_path, 'r') as zip_file:
                # Le contenu principal est dans content.xml
                if 'content.xml' in zip_file.namelist():
                    content_xml = zip_file.read('content.xml')
                    
                    # Parser le XML
                    root = ET.fromstring(content_xml)
                    
                    # Extraire le texte des éléments text:p (paragraphes)
                    text_parts = []
                    for elem in root.iter():
                        if elem.tag.endswith('}p'):  # Paragraphes
                            text = ''.join(elem.itertext()).strip()
                            if text:
                                text_parts.append(text)
                    
                    return "\n".join(text_parts) if text_parts else ""
            
            return ""
            
        except Exception as e:
            self.logger.error(f"Erreur extraction ODT {file_path}: {str(e)}")
            return ""

    def _extract_from_ods(self, file_path: Path) -> str:
        """Extrait le texte d'un fichier ODS (OpenDocument Spreadsheet)"""
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            
            # Les fichiers ODS sont des archives ZIP
            with zipfile.ZipFile(file_path, 'r') as zip_file:
                # Le contenu principal est dans content.xml
                if 'content.xml' in zip_file.namelist():
                    content_xml = zip_file.read('content.xml')
                    
                    # Parser le XML
                    root = ET.fromstring(content_xml)
                    
                    # Extraire le texte des cellules
                    text_parts = []
                    for elem in root.iter():
                        if elem.tag.endswith('}table-cell'):  # Cellules
                            text = ''.join(elem.itertext()).strip()
                            if text:
                                text_parts.append(text)
                    
                    return " | ".join(text_parts) if text_parts else ""
            
            return ""
            
        except Exception as e:
            self.logger.error(f"Erreur extraction ODS {file_path}: {str(e)}")
            return ""

    def _extract_from_odp(self, file_path: Path) -> str:
        """Extrait le texte d'un fichier ODP (OpenDocument Presentation)"""
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            
            # Les fichiers ODP sont des archives ZIP
            with zipfile.ZipFile(file_path, 'r') as zip_file:
                # Le contenu principal est dans content.xml
                if 'content.xml' in zip_file.namelist():
                    content_xml = zip_file.read('content.xml')
                    
                    # Parser le XML
                    root = ET.fromstring(content_xml)
                    
                    # Extraire le texte des éléments de présentation
                    text_parts = []
                    for elem in root.iter():
                        if elem.tag.endswith('}p'):  # Paragraphes
                            text = ''.join(elem.itertext()).strip()
                            if text:
                                text_parts.append(text)
                    
                    return "\n".join(text_parts) if text_parts else ""
            
            return ""
            
        except Exception as e:
            self.logger.error(f"Erreur extraction ODP {file_path}: {str(e)}")
            return ""

    def get_supported_formats(self) -> list[str]:
        """Retourne la liste des formats supportés"""
        return get_supported_formats_keys(self.supported_formats)

    def is_format_supported(self, extension: str) -> bool:
        """Vérifie si un format est supporté"""
        return is_format_supported_in_dict(extension, self.supported_formats) 