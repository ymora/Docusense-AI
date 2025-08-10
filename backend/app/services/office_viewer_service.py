"""
Office Viewer Service for DocuSense AI
Converts Office documents to HTML for direct browser viewing
"""

from pathlib import Path
import tempfile
import shutil
from typing import Optional, Dict, Any
from sqlalchemy.orm import Session
import asyncio
import logging

from .base_service import BaseService, log_service_operation
from ..core.types import ServiceResponse
from ..core.media_formats import get_supported_formats_keys, is_format_supported_in_dict


class OfficeViewerService(BaseService):
    """
    Service de conversion des documents Office pour visualisation web
    """

    def __init__(self, db: Session):
        super().__init__(db)
        self.supported_formats = {
            # Documents Word
            'docx': self._convert_docx_to_html,
            'doc': self._convert_doc_to_html,
            'odt': self._convert_odt_to_html,
            
            # Tableurs Excel
            'xlsx': self._convert_xlsx_to_html,
            'xls': self._convert_xls_to_html,
            'ods': self._convert_ods_to_html,
            
            # Présentations PowerPoint
            'pptx': self._convert_pptx_to_html,
            'ppt': self._convert_ppt_to_html,
            'odp': self._convert_odp_to_html,
        }
        
        # Cache temporaire pour les conversions
        self._conversion_cache = {}

    @log_service_operation("convert_to_html")
    async def convert_to_html(self, file_path: str) -> ServiceResponse:
        """
        Convertit un fichier Office en HTML pour visualisation
        
        Args:
            file_path: Chemin vers le fichier Office
            
        Returns:
            ServiceResponse: HTML généré ou erreur
        """
        try:
            html_content = await self.safe_execute_async("convert_to_html_logic", self._convert_to_html_logic, file_path)
            return {"success": True, "data": {"html": html_content, "file_path": file_path}}
        except Exception as e:
            return {"success": False, "error": f"Erreur lors de la conversion: {str(e)}"}

    async def _convert_to_html_logic(self, file_path: str) -> str:
        """Logic for converting file to HTML"""
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"Fichier non trouvé: {file_path}")

        extension = file_path.suffix.lower().lstrip('.')
        
        if extension not in self.supported_formats:
            raise ValueError(f"Format non supporté pour la visualisation: {extension}")

        # Vérifier le cache
        cache_key = f"{file_path}_{file_path.stat().st_mtime}"
        if cache_key in self._conversion_cache:
            return self._conversion_cache[cache_key]

        # Convertir le fichier
        converter_func = self.supported_formats[extension]
        html_content = await asyncio.to_thread(converter_func, file_path)
        
        # Mettre en cache
        self._conversion_cache[cache_key] = html_content
        
        return html_content

    def _convert_docx_to_html(self, file_path: Path) -> str:
        """Convertit un fichier DOCX en HTML"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            html_parts = ['<div class="office-document">']
            
            # Titre du document
            html_parts.append(f'<h1 class="document-title">{file_path.stem}</h1>')
            
            # Paragraphes
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    style = paragraph.style.name if paragraph.style else 'Normal'
                    html_parts.append(f'<p class="paragraph {style.lower()}">{paragraph.text}</p>')
            
            # Tableaux
            for table in doc.tables:
                html_parts.append('<table class="document-table">')
                for row in table.rows:
                    html_parts.append('<tr>')
                    for cell in row.cells:
                        html_parts.append(f'<td>{cell.text}</td>')
                    html_parts.append('</tr>')
                html_parts.append('</table>')
            
            html_parts.append('</div>')
            return '\n'.join(html_parts)
            
        except ImportError:
            return self._fallback_conversion(file_path, "DOCX")
        except Exception as e:
            self.logger.error(f"Erreur conversion DOCX: {e}")
            return self._fallback_conversion(file_path, "DOCX")

    def _convert_xlsx_to_html(self, file_path: Path) -> str:
        """Convertit un fichier XLSX en HTML"""
        try:
            import openpyxl
            
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            html_parts = ['<div class="office-spreadsheet">']
            
            # Titre du document
            html_parts.append(f'<h1 class="document-title">{file_path.stem}</h1>')
            
            # Parcourir toutes les feuilles
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                html_parts.append(f'<h2 class="sheet-title">{sheet_name}</h2>')
                html_parts.append('<table class="spreadsheet-table">')
                
                # Trouver les limites des données
                max_row = sheet.max_row
                max_col = sheet.max_column
                
                for row in range(1, min(max_row + 1, 100)):  # Limiter à 100 lignes
                    html_parts.append('<tr>')
                    for col in range(1, min(max_col + 1, 20)):  # Limiter à 20 colonnes
                        cell = sheet.cell(row=row, column=col)
                        value = cell.value if cell.value is not None else ''
                        html_parts.append(f'<td>{value}</td>')
                    html_parts.append('</tr>')
                
                html_parts.append('</table>')
            
            html_parts.append('</div>')
            return '\n'.join(html_parts)
            
        except ImportError:
            return self._fallback_conversion(file_path, "XLSX")
        except Exception as e:
            self.logger.error(f"Erreur conversion XLSX: {e}")
            return self._fallback_conversion(file_path, "XLSX")

    def _convert_pptx_to_html(self, file_path: Path) -> str:
        """Convertit un fichier PPTX en HTML"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            html_parts = ['<div class="office-presentation">']
            
            # Titre du document
            html_parts.append(f'<h1 class="document-title">{file_path.stem}</h1>')
            
            # Parcourir toutes les diapositives
            for i, slide in enumerate(prs.slides, 1):
                html_parts.append(f'<div class="slide" id="slide-{i}">')
                html_parts.append(f'<h2 class="slide-title">Diapositive {i}</h2>')
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        html_parts.append(f'<p class="slide-text">{shape.text}</p>')
                
                html_parts.append('</div>')
            
            html_parts.append('</div>')
            return '\n'.join(html_parts)
            
        except ImportError:
            return self._fallback_conversion(file_path, "PPTX")
        except Exception as e:
            self.logger.error(f"Erreur conversion PPTX: {e}")
            return self._fallback_conversion(file_path, "PPTX")

    def _convert_doc_to_html(self, file_path: Path) -> str:
        """Convertit un fichier DOC en HTML (via extraction de texte)"""
        try:
            # Utiliser le service d'extraction existant
            from .document_extractor_service import DocumentExtractorService
            extractor = DocumentExtractorService(self.db)
            
            # Extraire le texte
            text = extractor._extract_from_doc(file_path)
            
            # Convertir en HTML simple
            html_parts = ['<div class="office-document">']
            html_parts.append(f'<h1 class="document-title">{file_path.stem}</h1>')
            
            # Diviser le texte en paragraphes
            paragraphs = text.split('\n\n')
            for para in paragraphs:
                if para.strip():
                    html_parts.append(f'<p class="paragraph">{para.strip()}</p>')
            
            html_parts.append('</div>')
            return '\n'.join(html_parts)
            
        except Exception as e:
            self.logger.error(f"Erreur conversion DOC: {e}")
            return self._fallback_conversion(file_path, "DOC")

    def _convert_xls_to_html(self, file_path: Path) -> str:
        """Convertit un fichier XLS en HTML (via extraction de texte)"""
        try:
            # Utiliser le service d'extraction existant
            from .document_extractor_service import DocumentExtractorService
            extractor = DocumentExtractorService(self.db)
            
            # Extraire le texte
            text = extractor._extract_from_xls(file_path)
            
            # Convertir en HTML simple
            html_parts = ['<div class="office-spreadsheet">']
            html_parts.append(f'<h1 class="document-title">{file_path.stem}</h1>')
            html_parts.append('<div class="spreadsheet-content">')
            html_parts.append(f'<pre class="spreadsheet-text">{text}</pre>')
            html_parts.append('</div>')
            html_parts.append('</div>')
            
            return '\n'.join(html_parts)
            
        except Exception as e:
            self.logger.error(f"Erreur conversion XLS: {e}")
            return self._fallback_conversion(file_path, "XLS")

    def _convert_ppt_to_html(self, file_path: Path) -> str:
        """Convertit un fichier PPT en HTML (via extraction de texte)"""
        try:
            # Utiliser le service d'extraction existant
            from .document_extractor_service import DocumentExtractorService
            extractor = DocumentExtractorService(self.db)
            
            # Extraire le texte
            text = extractor._extract_from_ppt(file_path)
            
            # Convertir en HTML simple
            html_parts = ['<div class="office-presentation">']
            html_parts.append(f'<h1 class="document-title">{file_path.stem}</h1>')
            html_parts.append('<div class="presentation-content">')
            html_parts.append(f'<pre class="presentation-text">{text}</pre>')
            html_parts.append('</div>')
            html_parts.append('</div>')
            
            return '\n'.join(html_parts)
            
        except Exception as e:
            self.logger.error(f"Erreur conversion PPT: {e}")
            return self._fallback_conversion(file_path, "PPT")

    def _convert_odt_to_html(self, file_path: Path) -> str:
        """Convertit un fichier ODT en HTML"""
        try:
            # Utiliser le service d'extraction existant
            from .document_extractor_service import DocumentExtractorService
            extractor = DocumentExtractorService(self.db)
            
            # Extraire le texte
            text = extractor._extract_from_odt(file_path)
            
            # Convertir en HTML simple
            html_parts = ['<div class="office-document">']
            html_parts.append(f'<h1 class="document-title">{file_path.stem}</h1>')
            
            # Diviser le texte en paragraphes
            paragraphs = text.split('\n\n')
            for para in paragraphs:
                if para.strip():
                    html_parts.append(f'<p class="paragraph">{para.strip()}</p>')
            
            html_parts.append('</div>')
            return '\n'.join(html_parts)
            
        except Exception as e:
            self.logger.error(f"Erreur conversion ODT: {e}")
            return self._fallback_conversion(file_path, "ODT")

    def _convert_ods_to_html(self, file_path: Path) -> str:
        """Convertit un fichier ODS en HTML"""
        try:
            # Utiliser le service d'extraction existant
            from .document_extractor_service import DocumentExtractorService
            extractor = DocumentExtractorService(self.db)
            
            # Extraire le texte
            text = extractor._extract_from_ods(file_path)
            
            # Convertir en HTML simple
            html_parts = ['<div class="office-spreadsheet">']
            html_parts.append(f'<h1 class="document-title">{file_path.stem}</h1>')
            html_parts.append('<div class="spreadsheet-content">')
            html_parts.append(f'<pre class="spreadsheet-text">{text}</pre>')
            html_parts.append('</div>')
            html_parts.append('</div>')
            
            return '\n'.join(html_parts)
            
        except Exception as e:
            self.logger.error(f"Erreur conversion ODS: {e}")
            return self._fallback_conversion(file_path, "ODS")

    def _convert_odp_to_html(self, file_path: Path) -> str:
        """Convertit un fichier ODP en HTML"""
        try:
            # Utiliser le service d'extraction existant
            from .document_extractor_service import DocumentExtractorService
            extractor = DocumentExtractorService(self.db)
            
            # Extraire le texte
            text = extractor._extract_from_odp(file_path)
            
            # Convertir en HTML simple
            html_parts = ['<div class="office-presentation">']
            html_parts.append(f'<h1 class="document-title">{file_path.stem}</h1>')
            html_parts.append('<div class="presentation-content">')
            html_parts.append(f'<pre class="presentation-text">{text}</pre>')
            html_parts.append('</div>')
            html_parts.append('</div>')
            
            return '\n'.join(html_parts)
            
        except Exception as e:
            self.logger.error(f"Erreur conversion ODP: {e}")
            return self._fallback_conversion(file_path, "ODP")

    def _fallback_conversion(self, file_path: Path, format_type: str) -> str:
        """Conversion de fallback quand les bibliothèques ne sont pas disponibles"""
        return f"""
        <div class="office-document fallback">
            <h1 class="document-title">{file_path.stem}</h1>
            <div class="fallback-message">
                <p>⚠️ Visualisation limitée pour les fichiers {format_type}</p>
                <p>Le fichier <strong>{file_path.name}</strong> ne peut pas être affiché directement.</p>
                <p>Veuillez télécharger le fichier pour le consulter avec une application compatible.</p>
            </div>
        </div>
        """

    def get_supported_formats(self) -> list[str]:
        """Retourne la liste des formats supportés"""
        return get_supported_formats_keys(self.supported_formats)

    def is_format_supported(self, extension: str) -> bool:
        """Vérifie si un format est supporté"""
        return is_format_supported_in_dict(extension, self.supported_formats)
